"""Generate the 7-slide chairman presentation (PPTX).

Branding matches the existing chairman memo (AASTMT blue/navy).
Output: scripts/Carbon_Data_Trust_Platform_Chairman_Deck.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(OUTPUT_DIR, "Carbon_Data_Trust_Platform_Chairman_Deck.pptx")

# ── palette ─────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0x5C, 0xA8)
SKY = RGBColor(0xE8, 0xF1, 0xFA)
LIGHT = RGBColor(0xEE, 0xF4, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x55, 0x55, 0x55)
DARK_TEXT = RGBColor(0x22, 0x2A, 0x35)
GREEN = RGBColor(0x0E, 0x7C, 0x4F)
AMBER = RGBColor(0xB4, 0x6B, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Calibri"


# ── helpers ─────────────────────────────────────────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    return tb, tf


def set_run(run, text, size, color=DARK_TEXT, bold=False, italic=False):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_para(tf, text, size, color=DARK_TEXT, bold=False, first=False,
             space_before=6, space_after=0, level=0, align=None, bullet=False,
             italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    if align is not None:
        p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        _bullet(p, "•  ")
    run = p.add_run()
    set_run(run, text, size, color, bold=bold, italic=italic)
    return p


def _bullet(paragraph, marker):
    pPr = paragraph._p.get_or_add_pPr()
    buChar = pPr.makeelement(qn("a:buChar"), {"char": marker.replace("  ", "")})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": FONT})
    pPr.append(buChar)
    pPr.append(buFont)
    pPr.set("indent", "-228600")
    pPr.set("marL", "228600")


def header(slide, kicker, title, number=None):
    # top accent bar
    rect(slide, 0, 0, SLIDE_W, Inches(0.14), BLUE)
    # kicker
    tb, tf = textbox(slide, Inches(0.6), Inches(0.45), Inches(11.5), Inches(0.35))
    add_para(tf, kicker.upper(), 11, BLUE, bold=True, first=True, space_before=0)
    # title
    tb, tf = textbox(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.85))
    add_para(tf, title, 30, NAVY, bold=True, first=True, space_before=0)
    # rule under title
    rect(slide, Inches(0.6), Inches(1.62), Inches(2.2), Inches(0.04), BLUE)
    # slide number
    if number:
        tb, tf = textbox(slide, Inches(12.35), Inches(7.02), Inches(0.6), Inches(0.35))
        add_para(tf, f"{number}", 12, MUTED, first=True, space_before=0, align=PP_ALIGN.RIGHT)


def footer(slide, tagline):
    tb, tf = textbox(slide, Inches(0.6), Inches(7.05), Inches(9.0), Inches(0.35))
    add_para(tf, tagline, 9, MUTED, italic=True, first=True, space_before=0)


def pillar_chip(slide, x, y, w, h, text, fill):
    chip = rect(slide, x, y, w, h, fill)
    tf = chip.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, 14, WHITE, bold=True)
    return chip


def bullet_block(slide, x, y, w, h, items, size=15, gap=10, lead_color=BLUE):
    tb, tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
        else:
            lead, rest = None, it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": FONT})
        pPr.append(buChar)
        pPr.append(buFont)
        pPr.set("indent", "-205740")
        pPr.set("marL", "205740")
        if lead:
            rl = p.add_run()
            set_run(rl, lead + "  ", size, lead_color, bold=True)
        r = p.add_run()
        set_run(r, rest, size, DARK_TEXT)
    return tb


# ── slide builders ──────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # accent band
    rect(s, 0, Inches(3.02), SLIDE_W, Inches(0.06), BLUE)

    tb, tf = textbox(s, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.5))
    add_para(tf, "ARAB ACADEMY FOR SCIENCE, TECHNOLOGY & MARITIME TRANSPORT",
             14, SKY, bold=True, first=True, space_before=0)

    tb, tf = textbox(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.4))
    add_para(tf, "Carbon Data Trust Platform", 48, WHITE, bold=True, first=True, space_before=0)

    tb, tf = textbox(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.6))
    add_para(tf, "One place to trust your data, understand it, and act on it.",
             20, SKY, italic=True, first=True, space_before=0)

    # three pillars
    labels = ["TRUST", "INTELLIGENCE", "ACTION"]
    descs = ["Checked & traceable data", "AI that works with you", "Apps that drive results"]
    x = Inches(1.0)
    for lab, dsc in zip(labels, descs):
        pillar_chip(s, x, Inches(4.35), Inches(1.7), Inches(0.5), lab, BLUE)
        tb, tf = textbox(s, x, Inches(4.95), Inches(2.9), Inches(0.4))
        add_para(tf, dsc, 12, SKY, first=True, space_before=0)
        x += Inches(3.15)

    tb, tf = textbox(s, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5))
    add_para(tf, "Board of Trustees Briefing  ·  26 August 2026", 13, MUTED,
             first=True, space_before=0)
    return s


def slide_data_trust(prs):
    s = blank_slide(prs)
    header(s, "Dimension 1 · The Foundation", "The Data Trust Platform", 2)

    items = [
        ("A single trusted home", "for all university data — not another database, but the backbone everything sits on."),
        ("Every number is checked, scored, and traced", "back to its source, with a complete audit trail."),
        ("Who sees what is controlled precisely", "through capability-based access (CBAC) — the right eyes, nothing leaks."),
        ("Built to grow", "— new apps plug in without rebuilding the core."),
    ]
    bullet_block(s, Inches(0.6), Inches(2.0), Inches(7.4), Inches(4.6), items, size=16, gap=16)

    # right callout card
    card = rect(s, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.4), LIGHT)
    tb, tf = textbox(s, Inches(8.75), Inches(2.35), Inches(3.6), Inches(3.9))
    add_para(tf, "What makes it a platform", 14, BLUE, bold=True, first=True, space_before=0)
    for txt in [
        "Quality checks on every dataset",
        "Master data, kept consistent",
        "Full lineage & audit trail",
        "Governed, versioned datasets",
    ]:
        add_para(tf, txt, 13.5, DARK_TEXT, space_before=10, bullet=True)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_why_trust(prs):
    s = blank_slide(prs)
    header(s, "Dimension 1 · Why it matters", "Trust, for the people at the top", 3)

    # three statement cards
    cards = [
        ("One version of the truth",
         "No more arguing over whose spreadsheet is right. Leadership sees the same, agreed numbers."),
        ("Answers you can defend",
         "Every figure links back to its source — audit-ready when regulators or auditors ask."),
        ("Software that helps you run",
         "Trusted data feeds the AI layer, turning the platform from a record-keeper into a decision tool."),
    ]
    x = Inches(0.6)
    for title, body in cards:
        rect(s, x, Inches(2.0), Inches(3.95), Inches(2.6), LIGHT)
        rect(s, x, Inches(2.0), Inches(3.95), Inches(0.08), BLUE)
        tb, tf = textbox(s, x + Inches(0.3), Inches(2.3), Inches(3.35), Inches(2.1))
        add_para(tf, title, 16, NAVY, bold=True, first=True, space_before=0)
        add_para(tf, body, 13, DARK_TEXT, space_before=8)
        x += Inches(4.15)

    # bottom banner
    banner = rect(s, Inches(0.6), Inches(5.0), Inches(12.15), Inches(1.1), NAVY)
    tb, tf = textbox(s, Inches(1.0), Inches(5.2), Inches(11.4), Inches(0.8))
    add_para(tf, "The result: leadership gets trustworthy answers, faster — the foundation every decision stands on.",
             15, WHITE, bold=True, first=True, space_before=0, align=PP_ALIGN.CENTER)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_pulse(prs):
    s = blank_slide(prs)
    header(s, "Dimension 2 · The Intelligence", "Pulse — your AI coworker", 4)

    items = [
        ("Not a chatbot.", "A coworker that plans, reasons, and learns from every conversation."),
        ("Knows your data, your rules, and your people", "— and gets smarter the more it works."),
        ("Checks its own answers", "before giving them. Built-in truthfulness guards mean no made-up numbers."),
        ("Never oversteps.", "It only touches what it's allowed to (CBAC), and every action is recorded."),
    ]
    bullet_block(s, Inches(0.6), Inches(2.0), Inches(7.4), Inches(4.6), items, size=16, gap=16)

    card = rect(s, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.4), LIGHT)
    tb, tf = textbox(s, Inches(8.75), Inches(2.35), Inches(3.6), Inches(3.9))
    add_para(tf, "Pulse works the way you do", 14, BLUE, bold=True, first=True, space_before=0)
    for txt in [
        "Plans multi-step tasks",
        "Learns from your feedback",
        "Orchestrates other agents",
        "Explains how it reasoned",
    ]:
        add_para(tf, txt, 13.5, DARK_TEXT, space_before=10, bullet=True)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_pulse_nextgen(prs):
    s = blank_slide(prs)
    header(s, "Dimension 2 · Why it's different", "Next-generation software", 5)

    # three feature rows
    rows = [
        ("A separate, reusable brain",
         "The intelligence isn't locked inside one app. It plugs into any system — this platform is just its first home."),
        ("It grows safely",
         "New skills and tools are added on the outside. The core never breaks, so it gets stronger without risk."),
        ("It improves over time",
         "Every interaction feeds back — the software gets smarter with use, the way a good team does."),
    ]
    y = Inches(2.0)
    for title, body in rows:
        rect(s, Inches(0.6), y, Inches(0.14), Inches(1.15), BLUE)
        tb, tf = textbox(s, Inches(1.0), y + Inches(0.05), Inches(11.6), Inches(1.1))
        add_para(tf, title, 17, NAVY, bold=True, first=True, space_before=0)
        add_para(tf, body, 14, DARK_TEXT, space_before=4)
        y += Inches(1.5)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_carbon_footprint(prs):
    s = blank_slide(prs)
    header(s, "Dimension 3 · The Domain App", "Carbon Footprint — from coverage to action", 6)

    items = [
        ("Not data collection.", "A decision tool for emissions across the whole campus."),
        ("It answers:", "where do we emit, how much, where's the trend, and what should we do about it?"),
        ("From A to B:", "knowing what we emit (A) → actually reducing it, tracked against real targets (B)."),
        ("And more:", "audit-ready GHG inventory (Scope 1/2/3, SBTi targets), campus coverage, forecasting."),
    ]
    bullet_block(s, Inches(0.6), Inches(2.0), Inches(12.15), Inches(4.4), items, size=16, gap=16)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_carbon_answers(prs):
    s = blank_slide(prs)
    header(s, "Dimension 3 · The Domain App", "What Carbon Footprint answers", 7)

    cards = [
        ("Where do we emit?",
         "Broken down by campus, building, fleet, and activity — so you see the hotspots, not one blurry total."),
        ("How much?",
         "Real, measured tonnes — calculated from activity data, never a guess."),
        ("Where's the trend?",
         "Up, down, or flat — and why. You see the direction before it becomes a problem."),
        ("What do we do?",
         "Ranked actions with the expected impact of each — so effort goes where it counts."),
    ]
    positions = [
        (Inches(0.6), Inches(2.0)),
        (Inches(6.75), Inches(2.0)),
        (Inches(0.6), Inches(4.15)),
        (Inches(6.75), Inches(4.15)),
    ]
    for (title, body), (x, y) in zip(cards, positions):
        rect(s, x, y, Inches(5.95), Inches(1.95), LIGHT)
        rect(s, x, y, Inches(5.95), Inches(0.08), BLUE)
        tb, tf = textbox(s, x + Inches(0.3), y + Inches(0.25), Inches(5.35), Inches(1.55))
        add_para(tf, title, 16, NAVY, bold=True, first=True, space_before=0)
        add_para(tf, body, 12.5, DARK_TEXT, space_before=6)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_carbon_coverage(prs):
    s = blank_slide(prs)
    header(s, "Dimension 3 · The Domain App", "Coverage that's audit-ready", 8)

    cards = [
        ("Scope 1", "Direct emissions",
         "Fuel we burn ourselves — fleet, generators, campus equipment.", BLUE),
        ("Scope 2", "Purchased energy",
         "The electricity and cooling we buy — powering buildings and labs.", NAVY),
        ("Scope 3", "The wider chain",
         "Travel, procurement, waste — emissions from everything we buy and do.", GREEN),
    ]
    x = Inches(0.6)
    for title, sub, body, color in cards:
        rect(s, x, Inches(2.0), Inches(3.95), Inches(2.9), LIGHT)
        rect(s, x, Inches(2.0), Inches(3.95), Inches(0.08), color)
        tb, tf = textbox(s, x + Inches(0.3), Inches(2.25), Inches(3.35), Inches(2.5))
        add_para(tf, title, 20, color, bold=True, first=True, space_before=0)
        add_para(tf, sub, 13, MUTED, bold=True, space_before=2)
        add_para(tf, body, 12.5, DARK_TEXT, space_before=8)
        x += Inches(4.15)

    banner = rect(s, Inches(0.6), Inches(5.2), Inches(12.15), Inches(1.2), NAVY)
    tb, tf = textbox(s, Inches(1.0), Inches(5.35), Inches(11.4), Inches(0.9))
    add_para(tf, "Aligned to GHG Protocol · SBTi targets · full campus coverage",
             15, WHITE, bold=True, first=True, space_before=0, align=PP_ALIGN.CENTER)
    add_para(tf, "Auditors get a complete, traceable inventory — not a spreadsheet.",
             12.5, SKY, space_before=4, align=PP_ALIGN.CENTER)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_carbon_journey(prs):
    s = blank_slide(prs)
    header(s, "Dimension 3 · The Domain App", "From A to B — the journey", 9)

    # A card
    rect(s, Inches(0.6), Inches(2.0), Inches(5.95), Inches(2.6), LIGHT)
    rect(s, Inches(0.6), Inches(2.0), Inches(5.95), Inches(0.08), AMBER)
    tb, tf = textbox(s, Inches(0.95), Inches(2.25), Inches(5.3), Inches(2.2))
    add_para(tf, "A — Know what we emit", 18, NAVY, bold=True, first=True, space_before=0)
    for txt in [
        "Measure every source, campus-wide",
        "Verify the numbers — trusted & traceable",
    ]:
        add_para(tf, txt, 13.5, DARK_TEXT, space_before=8, bullet=True)

    # B card
    rect(s, Inches(6.85), Inches(2.0), Inches(5.95), Inches(2.6), LIGHT)
    rect(s, Inches(6.85), Inches(2.0), Inches(5.95), Inches(0.08), GREEN)
    tb, tf = textbox(s, Inches(7.2), Inches(2.25), Inches(5.3), Inches(2.2))
    add_para(tf, "B — Reduce against targets", 18, NAVY, bold=True, first=True, space_before=0)
    for txt in [
        "Set real reduction targets (SBTi)",
        "Track progress & forecast the outcome",
    ]:
        add_para(tf, txt, 13.5, DARK_TEXT, space_before=8, bullet=True)

    # journey steps strip
    steps = ["Measure", "Verify", "Set targets", "Act", "Forecast", "Report"]
    x = Inches(0.6)
    step_w = Inches(1.85)
    gap = Inches(0.21)
    y = Inches(5.1)
    for i, st in enumerate(steps):
        chip = rect(s, x, y, step_w, Inches(0.6), BLUE if i < 4 else NAVY)
        tf = chip.text_frame
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, st, 14, WHITE, bold=True)
        x += step_w + gap

    tb, tf = textbox(s, Inches(0.6), Inches(5.95), Inches(12.15), Inches(0.6))
    add_para(tf, "Forecasting: if we take action X, emissions drop to Y by 2030 — you decide with the end in sight.",
             13.5, MUTED, italic=True, first=True, space_before=0, align=PP_ALIGN.CENTER)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


def slide_big_picture(prs):
    s = blank_slide(prs)
    header(s, "One platform, three layers", "The big picture — and the ask", 10)

    # three layers stacked left
    layers = [
        ("ACTION", "Domain apps on the trusted core", SKY, NAVY),
        ("INTELLIGENCE", "Pulse — the AI coworker", WHITE, BLUE),
        ("TRUST", "The Data Trust Platform", LIGHT, NAVY),
    ]
    y = Inches(2.0)
    for lab, dsc, fill, txt in layers:
        chip = rect(s, Inches(0.6), y, Inches(2.3), Inches(0.7), fill)
        tf = chip.text_frame
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, lab, 14, txt, bold=True)
        tb, tf = textbox(s, Inches(3.1), y + Inches(0.08), Inches(4.6), Inches(0.55))
        add_para(tf, dsc, 14, DARK_TEXT, first=True, space_before=0)
        y += Inches(0.85)

    # right: AASTMT next apps
    rect(s, Inches(8.2), Inches(1.95), Inches(4.55), Inches(2.9), LIGHT)
    tb, tf = textbox(s, Inches(8.5), Inches(2.15), Inches(4.0), Inches(2.6))
    add_para(tf, "Already mapped for AASTMT", 13, BLUE, bold=True, first=True, space_before=0)
    for txt in [
        "Healthy Foods Factory (sales & operations AI)",
        "Facility Management",
        "Academic Portfolio",
        "HR & performance dashboards",
    ]:
        add_para(tf, txt, 12.5, DARK_TEXT, space_before=7, bullet=True)

    # the ask banner
    banner = rect(s, Inches(0.6), Inches(5.0), Inches(12.15), Inches(1.4), NAVY)
    tb, tf = textbox(s, Inches(1.0), Inches(5.15), Inches(11.4), Inches(1.1))
    add_para(tf, "The ask", 12, SKY, bold=True, first=True, space_before=0)
    add_para(tf, "Approve the roadmap and endorse the next phase of rollout.", 17,
             WHITE, bold=True, space_before=4)

    footer(s, "Carbon Data Trust Platform · AASTMT")
    return s


# ── build ───────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_data_trust(prs)
    slide_why_trust(prs)
    slide_pulse(prs)
    slide_pulse_nextgen(prs)
    slide_carbon_footprint(prs)
    slide_carbon_answers(prs)
    slide_carbon_coverage(prs)
    slide_carbon_journey(prs)
    slide_big_picture(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
