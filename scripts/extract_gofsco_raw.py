#!/usr/bin/env python
"""Extract all GOFSCO raw documents to text files, then dump all People-app
employees + metadata from PostgreSQL to structured JSON files.

Usage:
    python scripts/extract_gofsco_raw.py              # extract raw docs only (no Django)
    python scripts/extract_gofsco_raw.py --people     # also dump People app data (needs Django)
    python scripts/extract_gofsco_raw.py --all        # both
"""
import argparse
import json
import os
import sys
from pathlib import Path

RAW_DIR = Path("/home/ahmed/aast/carbon/raw/GOFSCO app/20260728")
OUT_DIR = Path("/home/ahmed/aast/carbon/raw/GOFSCO app/_extracted")
PEOPLE_OUT = Path("/home/ahmed/aast/carbon/raw/GOFSCO app/_extracted/people")
TEXT_EXTS = {".txt", ".md", ".csv", ".tsv"}
IMAGE_EXTS = {".jpeg", ".jpg", ".png", ".gif", ".webp", ".bmp", ".tiff"}


def extract_image_ocr(path: Path) -> str:
    """OCR a single image using tesseract (eng+ara). Returns empty string on failure."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return "[pytesseract not installed — run: pip install pytesseract]"
    try:
        img = Image.open(str(path))
        # Try Arabic + English; fall back to English-only if ara lang pack missing
        try:
            text = pytesseract.image_to_string(img, lang="ara+eng")
        except Exception:
            text = pytesseract.image_to_string(img, lang="eng")
        return text.strip()
    except Exception as e:  # noqa: BLE001
        return f"[OCR error: {e}]"


def extract_docx(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    parts = []
    for p in d.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"\n=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if txt.strip():
                    parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            txt = page.extract_text() or ""
            parts.append(f"\n=== Page {i} ===\n{txt}")
            for table in page.extract_tables():
                for row in table:
                    cells = ["" if c is None else str(c).strip() for c in row]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"\n=== Sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c).strip() for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_file(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in TEXT_EXTS:
        return extract_text(path)
    if ext in IMAGE_EXTS:
        return extract_image_ocr(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".xlsx":
        return extract_xlsx(path)
    return None


def extract_raw_docs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    manifest = []
    for f in sorted(RAW_DIR.rglob("*")):
        if not f.is_file():
            continue
        if f.suffix.lower() in {".ini"}:
            manifest.append({"file": str(f.relative_to(RAW_DIR)), "extracted": False,
                             "reason": "system file — skipped"})
            continue
        try:
            text = extract_file(f)
        except Exception as e:  # noqa: BLE001
            manifest.append({"file": str(f.relative_to(RAW_DIR)), "extracted": False,
                             "reason": f"extraction error: {e}"})
            continue
        if text is None:
            manifest.append({"file": str(f.relative_to(RAW_DIR)), "extracted": False,
                             "reason": "unsupported type"})
            continue
        # Flatten nested paths into a safe filename
        rel = f.relative_to(RAW_DIR)
        safe_name = str(rel).replace("/", "__").replace("\\", "__")
        out = OUT_DIR / f"{safe_name}.txt"
        out.write_text(text, encoding="utf-8")
        manifest.append({"file": str(rel), "extracted": True,
                         "output": str(out.relative_to(OUT_DIR.parent)),
                         "chars": len(text)})
    manifest.sort(key=lambda m: m["file"])
    (OUT_DIR / "_manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    ok = sum(1 for m in manifest if m.get("extracted"))
    print(f"Extracted {ok}/{len(manifest)} files -> {OUT_DIR}")
    return manifest


# ────────────────────────────────────────────────────────────────────────
# People app dump
# ────────────────────────────────────────────────────────────────────────

def _dec(v):
    return None if v is None else str(v)


def _date(v):
    return None if v is None else v.isoformat()


def _dt(v):
    return None if v is None else v.isoformat()


def dump_people():
    # Ensure backend/ is on sys.path so `config.settings` resolves regardless of cwd.
    BACKEND = Path("/home/ahmed/aast/carbon/backend")
    if str(BACKEND) not in sys.path:
        sys.path.insert(0, str(BACKEND))
    os.environ.setdefault("DJANGO_SETTINGS_MODULE", "config.settings")
    import django
    django.setup()

    from people import models as p
    from mdm import models as m

    PEOPLE_OUT.mkdir(parents=True, exist_ok=True)

    def w(name, obj):
        (PEOPLE_OUT / f"{name}.json").write_text(
            json.dumps(obj, indent=2, ensure_ascii=False, default=str), encoding="utf-8")

    # ── Org units ──
    ous = [{
        "id": o.id, "code": o.code, "name": o.name, "slug": getattr(o, "slug", ""),
        "parent_id": o.parent_id, "is_active": getattr(o, "is_active", True),
        "lifecycle_state": getattr(o, "lifecycle_state", ""),
    } for o in m.OrgUnit.objects.all().order_by("id")]
    w("org_units", ous)

    # ── Reference sets / values (governed enums) ──
    refsets = []
    for rs in m.ReferenceSet.objects.all().order_by("id"):
        vals = [{"code": v.code, "label": v.label, "is_active": v.is_active,
                 "sort_order": getattr(v, "sort_order", 0)}
                for v in rs.values.all().order_by("sort_order", "code")]
        refsets.append({"id": rs.id, "name": rs.name, "slug": rs.slug,
                        "description": rs.description, "is_active": rs.is_active,
                        "lifecycle_state": getattr(rs, "lifecycle_state", ""),
                        "values": vals})
    w("reference_sets", refsets)

    # ── Positions ──
    positions = [{
        "id": x.id, "org_unit_id": x.org_unit_id, "code": x.code, "title": x.title,
        "grade": x.grade, "reports_to_id": x.reports_to_id,
        "is_management": x.is_management, "status": x.status,
        "fte": _dec(x.fte), "job_family_code": x.job_family_code,
    } for x in p.Position.objects.all().order_by("id")]
    w("positions", positions)

    # ── Compensation components ──
    comps = [{
        "id": x.id, "code": x.code, "name": x.name, "name_ar": x.name_ar,
        "direction": x.direction, "category": x.category,
        "is_eosi_base": x.is_eosi_base, "is_gosi_base": x.is_gosi_base,
        "is_wps_relevant": x.is_wps_relevant, "is_taxable": x.is_taxable,
        "is_variable": x.is_variable, "sort_order": x.sort_order,
        "valid_from": _date(x.valid_from), "valid_to": _date(x.valid_to),
        "is_active": x.is_active,
    } for x in p.CompensationComponent.objects.all().order_by("direction", "sort_order", "code")]
    w("compensation_components", comps)

    # ── Benefit types ──
    bts = [{
        "id": x.id, "code": x.code, "name": x.name, "category": x.category,
        "is_eosi_base": x.is_eosi_base, "is_taxable": x.is_taxable,
    } for x in p.BenefitType.objects.all().order_by("category", "code")]
    w("benefit_types", bts)

    # ── Compliance rules ──
    rules = [{
        "id": x.id, "rule_id": x.rule_id, "version": x.version, "name": x.name,
        "description": x.description, "jurisdiction": x.jurisdiction,
        "category": x.category, "effective_date": _date(x.effective_date),
        "formula_ref": x.formula_ref, "source_citation": x.source_citation,
        "inputs_schema": x.inputs_schema, "is_authoritative": x.is_authoritative,
        "provenance": x.provenance, "test_cases": x.test_cases,
    } for x in p.ComplianceRule.objects.all().order_by("category", "rule_id")]
    w("compliance_rules", rules)

    # ── Employees (the core) ──
    employees = []
    for e in p.Employee.objects.all().order_by("employee_no"):
        employees.append({
            "id": e.id, "employee_no": e.employee_no, "full_name": e.full_name,
            "name_en_given": e.name_en_given, "name_en_family": e.name_en_family,
            "name_ar_given": e.name_ar_given, "name_ar_family": e.name_ar_family,
            "civil_id": e.civil_id, "date_of_birth": _date(e.date_of_birth),
            "gender": e.gender, "nationality": e.nationality,
            "nationality_code": e.nationality_code,
            "employment_type_code": e.employment_type_code,
            "contract_type_code": e.contract_type_code,
            "kuwaitization": e.kuwaitization, "manager_id": e.manager_id,
            "position_id": e.position_id, "org_unit_id": e.org_unit_id,
            "basic_salary": _dec(e.basic_salary), "join_date": _date(e.join_date),
            "rotation": e.rotation, "is_active": e.is_active,
            "created_at": _dt(e.created_at), "updated_at": _dt(e.updated_at),
            "certifications": [
                {"cert_type": c.cert_type, "number": c.number,
                 "issued_date": _date(c.issued_date), "expiry_date": _date(c.expiry_date),
                 "notes": c.notes}
                for c in e.certifications.all().order_by("expiry_date")
            ],
            "benefits": [
                {"benefit_type_code": b.benefit_type.code, "benefit_type_name": b.benefit_type.name,
                 "monthly_amount": _dec(b.monthly_amount),
                 "effective_start": _date(b.effective_start),
                 "effective_end": _date(b.effective_end)}
                for b in e.benefits.all().order_by("benefit_type__code")
            ],
            "leave_entitlements": [
                {"year": l.year, "leave_type": l.leave_type,
                 "entitled_days": _dec(l.entitled_days), "used_days": _dec(l.used_days),
                 "carried_forward": _dec(l.carried_forward), "notes": l.notes}
                for l in e.leave_entitlements.all().order_by("year", "leave_type")
            ],
            "loans": [
                {"loan_type": ln.loan_type, "principal": _dec(ln.principal),
                 "interest_rate": _dec(ln.interest_rate), "term_months": ln.term_months,
                 "start_date": _date(ln.start_date), "status": ln.status, "notes": ln.notes}
                for ln in e.loans.all().order_by("-start_date")
            ],
            "compensation_lines": [
                {"component_code": cl.component.code, "component_name": cl.component.name,
                 "amount": _dec(cl.amount), "currency": cl.currency, "frequency": cl.frequency,
                 "effective_start": _date(cl.effective_start), "effective_end": _date(cl.effective_end),
                 "reason_note": cl.reason_note, "is_verified": cl.is_verified,
                 "verified_at": _dt(cl.verified_at)}
                for cl in e.compensation_lines.all().order_by("component__code", "-effective_start")
            ],
            "rotation_schedules": [
                {"pattern": r.pattern, "start_date": _date(r.start_date),
                 "config": r.config, "is_active": r.is_active}
                for r in e.rotation_schedules.all().order_by("-start_date")
            ],
            "attendance": [
                {"date": _date(a.date), "hours_worked": _dec(a.hours_worked),
                 "overtime_hours": _dec(a.overtime_hours), "status": a.status}
                for a in e.attendance.all().order_by("-date")
            ],
            "permissions": [
                {"date": _date(ap.date), "permission_type": ap.permission_type,
                 "hours": _dec(ap.hours), "approved": ap.approved, "notes": ap.notes}
                for ap in e.permissions.all().order_by("-date")
            ],
            "leave_records": [
                {"leave_type": lr.leave_type, "start_date": _date(lr.start_date),
                 "end_date": _date(lr.end_date), "days": _dec(lr.days),
                 "status": lr.status, "calendar_split": lr.calendar_split}
                for lr in e.leave_records.all().order_by("-start_date")
            ],
            "payslip_lines": [
                {"payroll_run_id": pl.payroll_run_id, "line_type": pl.line_type,
                 "amount": _dec(pl.amount), "rule_id": pl.rule_id,
                 "rule_version": pl.rule_version, "inputs": pl.inputs}
                for pl in e.payslip_lines.all().order_by("id")
            ],
        })
    w("employees", employees)

    # ── Personnel events ──
    events = [{
        "id": x.id, "entity_type": x.entity_type, "entity_id": x.entity_id,
        "event_kind": x.event_kind, "effective_date": _date(x.effective_date),
        "recorded_at": _dt(x.recorded_at), "recorded_by_id": x.recorded_by_id,
        "before": x.before, "after": x.after, "notes": x.notes,
    } for x in p.PersonnelEvent.objects.all().order_by("-effective_date", "-recorded_at")]
    w("personnel_events", events)

    # ── Compensation plans ──
    plans = [{
        "id": x.id, "org_unit_id": x.org_unit_id, "pay_grade_code": x.pay_grade_code,
        "job_family_code": x.job_family_code, "component_code": x.component.code,
        "amount": _dec(x.amount), "currency": x.currency, "frequency": x.frequency,
        "effective_start": _date(x.effective_start), "effective_end": _date(x.effective_end),
        "is_active": x.is_active,
    } for x in p.CompensationPlan.objects.all().order_by("pay_grade_code", "component__code")]
    w("compensation_plans", plans)

    # ── Payroll runs ──
    runs = [{
        "id": x.id, "org_unit_id": x.org_unit_id, "period_start": _date(x.period_start),
        "period_end": _date(x.period_end), "status": x.status,
        "created_at": _dt(x.created_at), "committed_at": _dt(x.committed_at),
    } for x in p.PayrollRun.objects.all().order_by("-period_start")]
    w("payroll_runs", runs)

    summary = {
        "employees": len(employees),
        "org_units": len(ous),
        "positions": len(positions),
        "compensation_components": len(comps),
        "benefit_types": len(bts),
        "compliance_rules": len(rules),
        "reference_sets": len(refsets),
        "personnel_events": len(events),
        "compensation_plans": len(plans),
        "payroll_runs": len(runs),
    }
    w("_summary", summary)
    print(f"People dump -> {PEOPLE_OUT}")
    print(json.dumps(summary, indent=2))
    return summary


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--people", action="store_true")
    ap.add_argument("--all", action="store_true")
    args = ap.parse_args()
    do_people = args.people or args.all
    do_raw = not args.people or args.all
    if do_raw:
        extract_raw_docs()
    if do_people:
        dump_people()


if __name__ == "__main__":
    main()
